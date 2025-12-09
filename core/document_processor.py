"""
Document Processor - Extract text from various file formats
"""

import os
from pathlib import Path
import fitz  # PyMuPDF
from pptx import Presentation
from docx import Document
import docx2txt

class DocumentProcessor:
    def __init__(self):
        self.supported_formats = {'.pdf', '.docx', '.pptx', '.txt'}
    
    async def extract_text(self, file_path: str) -> str:
        """Extract text from various document formats"""
        try:
            file_path = Path(file_path)
            extension = file_path.suffix.lower()
            
            if extension not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {extension}")
            
            if extension == '.pdf':
                return await self._extract_from_pdf(file_path)
            elif extension == '.docx':
                return await self._extract_from_docx(file_path)
            elif extension == '.pptx':
                return await self._extract_from_pptx(file_path)
            elif extension == '.txt':
                return await self._extract_from_txt(file_path)
            
        except Exception as e:
            print(f"Error extracting text from {file_path}: {e}")
            return ""
    
    async def _extract_from_pdf(self, file_path: Path) -> str:
        """Extract text from PDF using PyMuPDF"""
        try:
            doc = fitz.open(file_path)
            pages = []
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                pages.append(page.get_text())
                pages.append("\n\n")  # Add page separator
            
            doc.close()
            return "".join(pages).strip()
            
        except Exception as e:
            print(f"Error extracting PDF: {e}")
            return ""
    
    async def _extract_from_docx(self, file_path: Path) -> str:
        """Extract text from DOCX"""
        try:
            # Method 1: Using python-docx
            doc = Document(file_path)
            parts = []
            
            for paragraph in doc.paragraphs:
                parts.append(paragraph.text)
                parts.append("\n")
            
            # Extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
                        parts.append(" ")
                    parts.append("\n")
            
            return "".join(parts).strip()
            
        except Exception as e:
            try:
                # Method 2: Using docx2txt as fallback
                text = docx2txt.process(str(file_path))
                return text.strip()
            except Exception as e2:
                print(f"Error extracting DOCX: {e}, {e2}")
    async def _extract_from_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentations"""
        try:
            prs = Presentation(file_path)
            parts = []
            slide_count = 0
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_count += 1
                parts.append(f"Slide {slide_num}:\n")
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        parts.append(shape.text.strip())
                        parts.append("\n")
                
                parts.append("\n")
            
            text = "".join(parts).strip()
            print(f"📊 PowerPoint processed: {slide_count} slides, {len(text)} characters extracted")
            return text
            
        except Exception as e:
            print(f"Error extracting PPTX: {e}")
            return ""ion as e:
            print(f"Error extracting PPTX: {e}")
            return ""
    
    async def _extract_from_txt(self, file_path: Path) -> str:
        """Extract text from plain text files"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                return file.read()
        except UnicodeDecodeError:
            # Try with different encoding
            try:
                with open(file_path, 'r', encoding='latin-1') as file:
                    return file.read()
            except Exception as e:
                print(f"Error reading text file: {e}")
                return ""
        except Exception as e:
            print(f"Error extracting TXT: {e}")
            return ""
    
    def is_supported_format(self, file_path: str) -> bool:
        """Check if file format is supported"""
        extension = Path(file_path).suffix.lower()
        return extension in self.supported_formats
    
    async def extract_text_from_bytes(self, file_content: bytes, filename: str) -> str:
        """Extract text from file content bytes"""
        try:
            # Create a temporary file to process the content
            import tempfile
            
            file_extension = Path(filename).suffix.lower()
            
            if file_extension not in self.supported_formats:
                raise ValueError(f"Unsupported file format: {file_extension}")
            
            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix=file_extension, delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            try:
                # Extract text using the existing method
                text = await self.extract_text(temp_file_path)
                return text
            finally:
                # Clean up temporary file
                try:
                    os.unlink(temp_file_path)
                except:
                    pass
                    
        except Exception as e:
            print(f"Error extracting text from bytes: {e}")
            return ""
    
    def get_file_info(self, file_path: str) -> dict:
        """Get basic information about the file"""
        try:
            path = Path(file_path)
            return {
                "name": path.name,
                "extension": path.suffix.lower(),
                "size": path.stat().st_size,
                "supported": self.is_supported_format(file_path)
            }
        except Exception as e:
            return {"error": str(e)}